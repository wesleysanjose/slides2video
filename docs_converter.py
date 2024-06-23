
from abc import ABC, abstractmethod
import logging
import os
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
import magic
import aspose.slides as slides
import subprocess

# Define an abstract class for converting documents to JPEG

# Set up logging
# Configure logging to include filename, class name, and method name
logging.basicConfig(level=logging.INFO,
                    format='%(asctime)s - [%(filename)s - %(name)s - %(funcName)s] - %(levelname)s - %(message)s',
                    handlers=[
                        logging.FileHandler("document_conversion.log"),
                        logging.StreamHandler()
                    ])


class ConverterFactory:
    @staticmethod
    def create_converter_simple(file_path, images_dir):
        _, ext = os.path.splitext(file_path)
        if ext.lower() in ['.pdf']:
            return PDFToJPEGConverter(file_path, images_dir)
        elif ext.lower() in ['.ppt', '.pptx']:
            return PPTToJPEGConverter(file_path, images_dir)
        else:
            raise ValueError("Unsupported file type")

    def create_converter(file_path, images_dir):
        """
            Determines if the given file is a PowerPoint or PDF by inspecting the file content.

            Args:
            file_path (str): The path to the file.

            Returns:
            str: A message indicating the file type or if it's neither a PDF nor a PowerPoint file.
        """
        mime = magic.Magic(mime=True)
        mime_type = mime.from_file(file_path)

        if 'application/pdf' in mime_type:
            return PDFToJPEGConverter(file_path, images_dir)
        elif 'application/vnd.ms-powerpoint' in mime_type or 'application/vnd.openxmlformats-officedocument.presentationml.presentation' in mime_type:
            return PPTToJPEGConverter(file_path, images_dir)
        else:
            return None


class DocToJPEGConverter(ABC):
    def __init__(self, file_path, images_dir):
        self.file_path = file_path
        self.images_dir = images_dir

        # Set up logging
        logging.basicConfig(level=logging.INFO,
                            format='%(asctime)s - %(levelname)s - %(message)s')

        # Create the images directory if it doesn't exist
        if not os.path.exists(self.images_dir):
            os.makedirs(self.images_dir)
        else:
            # Clean the directory if it already exists
            self.clean_files_in_directory(self.images_dir)

    def clean_files_in_directory(self, directory):
        for filename in os.listdir(directory):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path):
                os.remove(file_path)
                logging.info(f"Deleted {file_path}")

    @abstractmethod
    def convert(self):
        """
        Abstract method to be implemented by subclasses to convert documents to JPEG.
        """
        pass


class PDFToJPEGConverter(DocToJPEGConverter):

    def __init__(self, file_path, images_dir):
        super().__init__(file_path, images_dir)

    def convert(self):
        """
        Convert each page of the PDF to a separate JPEG file.
        """
        try:
            # Convert PDF to a list of images
            images = convert_from_path(self.file_path)
            logging.info(
                f"Successfully converted PDF to images: {self.file_path}")

            # Save each image as JPEG
            for index, image in enumerate(images):
                image_path = f"{self.images_dir}/page_{index + 1}.jpg"
                image.save(image_path, 'JPEG')
                logging.info(f"Saved JPEG file at {image_path}")
        except FileNotFoundError:
            logging.error(
                "The specified PDF file was not found. Please check the file path.")
        except Exception as e:
            logging.error(f"An error occurred during conversion: {str(e)}")


class PPTToJPEGConverter:
    def __init__(self, ppt_file_path, output_dir):
        self.ppt_file_path = ppt_file_path
        self.output_dir = output_dir
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

    def convert(self):
        # Convert PPT to images using libreoffice
        command = [
            'libreoffice', '--headless', '--convert-to', 'jpg',
            '--outdir', self.output_dir, self.ppt_file_path
        ]
        subprocess.run(command, check=True)

        # Rename the output files to match the desired format
        for i, filename in enumerate(sorted(os.listdir(self.output_dir))):
            if filename.endswith('.jpg'):
                old_path = os.path.join(self.output_dir, filename)
                new_path = os.path.join(self.output_dir, f"slide_{i + 1}.jpg")
                os.rename(old_path, new_path)
                print(f"Saved slide {i + 1} as JPEG at '{new_path}'")

        # Extract notes using python-pptx
        prs = Presentation(self.ppt_file_path)
        for i, slide in enumerate(prs.slides):
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                notes_text = slide.notes_slide.notes_text_frame.text
                text_file = os.path.join(
                    self.output_dir, f"slide_{i + 1}_notes.txt")
                with open(text_file, 'w', encoding='utf-8') as f:
                    f.write(notes_text)
                print(f"Saved notes for slide {i + 1} at '{text_file}'")


# class PPTToJPEGConverter:
#     def __init__(self, ppt_file_path, output_dir):
#         self.ppt_file_path = ppt_file_path
#         self.output_dir = output_dir
#         if not os.path.exists(self.output_dir):
#             os.makedirs(self.output_dir)

#     def convert(self):
#         # Load the presentation
#         with slides.Presentation(self.ppt_file_path) as pres:
#             for i, slide in enumerate(pres.slides):
#                 # Define path for saving the slide image
#                 image_file = os.path.join(
#                     self.output_dir, f"slide_{i + 1}.jpg")
#                 # Save the slide as a JPEG image
#                 slide.get_thumbnail(1.0, 1.0).save(
#                     image_file, SaveFormat.JPEG)
#                 print(f"Saved slide {i + 1} as JPEG at '{image_file}'")

#                 # Export notes to a text file if they exist
#                 if slide.notes_slide_manager and slide.notes_slide_manager.notes_slide:
#                     notes_slide = slide.notes_slide_manager.notes_slide
#                     if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
#                         notes_text = notes_slide.notes_text_frame.text
#                         text_file = os.path.join(
#                             self.output_dir, f"slide_{i + 1}_notes.txt")
#                         with open(text_file, 'w') as f:
#                             f.write(notes_text)
#                         print(
#                             f"Saved notes for slide {i + 1} at '{text_file}'")
